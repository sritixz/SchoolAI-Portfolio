"""
Celery tasks for background presentation generation.

Each task:
  1. Calls the LLM once per slide for content.
  2. Optionally calls an image API for the visual_prompt.
  3. Builds a .pptx with python-pptx.
  4. Uploads the file to S3 and returns the download URL.
"""
import sys, os, json, asyncio, io, uuid, logging
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

import httpx
from celery import Task
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from worker.celery_app import celery
from config import settings
from services.s3 import upload_file

log = logging.getLogger(__name__)

# ── helpers ──────────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "695be6"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run(coro):
    """Run an async coroutine from sync Celery context."""
    return asyncio.get_event_loop().run_until_complete(coro)


async def _llm_slide(slide_number: int, total: int, topic: str, subject: str,
                     grade: str, board: str, tone: str, content_depth: str,
                     purpose: str, target_audience: str) -> dict:
    """Ask the LLM to generate a single slide's content."""
    from services.llm import chat_completion

    depth_rule = (
        "Max 3 bullets, high-level summary only."
        if content_depth == "Concise"
        else "Comprehensive explanation with data points and context."
    )
    tone_rule = {
        "Engaging":   "Use analogies, 'Did you know?' facts, and real-world examples.",
        "Formal":     "Professional, structured, academic language.",
        "Reflection": "Reflective prompts and deeper thinking questions.",
    }.get(tone, "")

    slide_types = ["title", "hook", "content", "content", "example",
                   "activity", "summary", "assessment"]
    slide_type = slide_types[min(slide_number - 1, len(slide_types) - 1)]
    if slide_number == total:
        slide_type = "assessment"

    prompt = f"""Generate slide {slide_number} of {total} for a {board} {subject} presentation on "{topic}" for {grade} ({target_audience}).
Slide type: {slide_type}. Purpose: {purpose}.
Content depth rule: {depth_rule}
Tone rule: {tone_rule}

Return ONLY valid JSON:
{{
  "number": {slide_number},
  "type": "{slide_type}",
  "title": "Slide title",
  "content": {{
    "bullets": ["bullet 1", "bullet 2"],
    "steps": [],
    "visual_prompt": "Vibrant, modern AI image description for this slide"
  }},
  "speaker_notes": "Teacher script for this slide",
  "engagement_prompt": "Question or activity for students",
  "vibrant_accent_color": "#695be6"
}}"""

    raw = await chat_completion([{"role": "user", "content": prompt}])
    raw = raw.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
    return json.loads(raw)


async def _generate_image(visual_prompt: str) -> bytes | None:
    """Call image API and return raw PNG bytes, or None if not configured."""
    if not settings.IMAGE_API_KEY:
        return None
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                settings.IMAGE_API_URL,
                headers={"Authorization": f"Bearer {settings.IMAGE_API_KEY}"},
                json={"prompt": visual_prompt, "n": 1, "size": "512x512",
                      "response_format": "url"},
            )
            resp.raise_for_status()
            img_url = resp.json()["data"][0]["url"]
            img_resp = await client.get(img_url)
            img_resp.raise_for_status()
            return img_resp.content
    except Exception as exc:
        log.warning("Image generation failed: %s", exc)
        return None


def _add_slide(prs: Presentation, slide_data: dict, img_bytes: bytes | None):
    """Append one slide to the pptx Presentation object."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)

    accent = _hex_to_rgb(slide_data.get("vibrant_accent_color", "#695be6"))
    W, H = prs.slide_width, prs.slide_height

    # ── accent bar at top ────────────────────────────────────────────────────
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # ── slide number badge ───────────────────────────────────────────────────
    badge = slide.shapes.add_shape(1, Inches(0.2), Inches(0.1), Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = str(slide_data.get("number", ""))
    tf.paragraphs[0].runs[0].font.size = Pt(9)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = accent
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── slide type label ─────────────────────────────────────────────────────
    type_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.1), Inches(2), Inches(0.35))
    tf = type_box.text_frame
    tf.text = (slide_data.get("type", "content")).upper()
    tf.paragraphs[0].runs[0].font.size = Pt(8)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

    # ── title ────────────────────────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), W - Inches(0.8), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = slide_data.get("title", "")
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(22)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(30, 30, 30)

    # ── image (right half) if available ─────────────────────────────────────
    content_right = W
    if img_bytes:
        try:
            img_stream = io.BytesIO(img_bytes)
            pic_left = W - Inches(3.8)
            slide.shapes.add_picture(img_stream, pic_left, Inches(1.5),
                                     Inches(3.6), Inches(2.8))
            content_right = pic_left - Inches(0.1)
        except Exception:
            pass

    # ── bullets / steps ──────────────────────────────────────────────────────
    content = slide_data.get("content", {})
    bullets = content.get("bullets", []) if isinstance(content, dict) else []
    steps   = content.get("steps",   []) if isinstance(content, dict) else []
    items   = bullets or [f"{i+1}. {s}" for i, s in enumerate(steps)]

    if items:
        body_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.55), content_right - Inches(0.5), Inches(3.2))
        tf = body_box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.runs[0].font.size = Pt(13)
            p.runs[0].font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(4)

    # ── speaker notes (amber box) ────────────────────────────────────────────
    notes = slide_data.get("speaker_notes", "")
    if notes:
        note_box = slide.shapes.add_textbox(
            Inches(0.4), H - Inches(1.4), W - Inches(0.8), Inches(0.9))
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = RGBColor(255, 251, 235)
        tf = note_box.text_frame
        tf.word_wrap = True
        tf.text = f"📋 {notes}"
        tf.paragraphs[0].runs[0].font.size = Pt(8)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(120, 80, 0)

    # ── engagement prompt (green box) ────────────────────────────────────────
    prompt_text = slide_data.get("engagement_prompt", "")
    if prompt_text:
        ep_box = slide.shapes.add_textbox(
            Inches(0.4), H - Inches(2.4), W - Inches(0.8), Inches(0.7))
        ep_box.fill.solid()
        ep_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
        tf = ep_box.text_frame
        tf.word_wrap = True
        tf.text = f"💬 {prompt_text}"
        tf.paragraphs[0].runs[0].font.size = Pt(8)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(20, 100, 50)


# ── Celery task ───────────────────────────────────────────────────────────────

@celery.task(bind=True, name="worker.tasks.generate_presentation")
def generate_presentation(self: Task, payload: dict) -> dict:
    """
    Build a full .pptx presentation slide-by-slide.

    payload keys: topic, subject, grade, board, tone, content_depth,
                  purpose, target_audience, num_slides, duration_minutes,
                  learning_objective, special_instructions, include_mini_quiz
    """
    topic            = payload["topic"]
    subject          = payload["subject"]
    grade            = payload.get("grade", "")
    board            = payload.get("board", "CBSE")
    tone             = payload.get("tone", "Engaging")
    content_depth    = payload.get("content_depth", "Concise")
    purpose          = payload.get("purpose", "teaching")
    target_audience  = payload.get("target_audience", "students")
    total            = int(payload.get("num_slides", 10))
    include_quiz     = payload.get("include_mini_quiz", False)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slides_data = []

    for i in range(1, total + 1):
        # ── update progress ──────────────────────────────────────────────────
        self.update_state(
            state="PROGRESS",
            meta={"current": i - 1, "total": total, "status": f"Generating slide {i} of {total}…"},
        )

        # ── LLM call for slide content ───────────────────────────────────────
        try:
            slide_data = _run(_llm_slide(
                i, total, topic, subject, grade, board,
                tone, content_depth, purpose, target_audience,
            ))
        except Exception as exc:
            log.error("LLM failed for slide %d: %s", i, exc)
            slide_data = {
                "number": i, "type": "content", "title": f"Slide {i}",
                "content": {"bullets": [], "steps": [], "visual_prompt": ""},
                "speaker_notes": "", "engagement_prompt": "",
                "vibrant_accent_color": "#695be6",
            }

        # ── image generation ─────────────────────────────────────────────────
        visual_prompt = (slide_data.get("content") or {}).get("visual_prompt", "")
        img_bytes = _run(_generate_image(visual_prompt)) if visual_prompt else None

        # ── build pptx slide ─────────────────────────────────────────────────
        _add_slide(prs, slide_data, img_bytes)
        slides_data.append(slide_data)

        self.update_state(
            state="PROGRESS",
            meta={"current": i, "total": total, "status": f"Slide {i} of {total} done"},
        )

    # ── upload to S3 ─────────────────────────────────────────────────────────
    self.update_state(state="PROGRESS",
                      meta={"current": total, "total": total, "status": "Uploading file…"})

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"{topic.replace(' ', '_')}_{uuid.uuid4().hex[:8]}.pptx"
    download_url = upload_file(buf.read(), filename, folder="presentations")

    return {
        "status": "DONE",
        "download_url": download_url,
        "slides": slides_data,
        "title": topic,
        "subject": subject,
        "grade": grade,
        "total_slides": total,
        "duration_minutes": payload.get("duration_minutes", 45),
    }
